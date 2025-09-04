import ollama    
import fitz  # PyMuPDF for PDFs
import docx  # python-docx for Word files
import pptx  # python-pptx for PowerPoint files
import os


def summarize_text(text, style="concise", length="short", character="Cartoon Character"):
    
    prompt = f"Summarize the following text in a {length} and {style} manner, in the style of {character}:\n\n{text}"

    try:
        response = ollama.generate(model="llama3.2:1b", prompt=prompt)
        return response["response"].strip()  
    except Exception as e:
        return f"Error: {str(e)}"
    


def answer_question(text):
    try:
        response = ollama.generate(model="llama3.2:1b", prompt=prompt)
        return response["response"].strip()  
    except Exception as e:
        return f"Error: {str(e)}"



def extract_text_from_file(file_path):
    """
    Extracts text from a given file (PDF, DOCX, or PPTX).

    Args:
        file_path (str): Path to the file.

    Returns:
        str: Extracted text.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    file_ext = os.path.splitext(file_path)[1].lower()

    if file_ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif file_ext == ".docx":
        return extract_text_from_docx(file_path)
    elif file_ext == ".pptx":
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError("Unsupported file type. Please provide a PDF, DOCX, or PPTX file.")

def extract_text_from_pdf(pdf_path):
    """Extracts text from a PDF file."""
    doc = fitz.open(pdf_path)
    text = "\n".join(page.get_text("text") for page in doc)
    return text

def extract_text_from_docx(docx_path):
    """Extracts text from a Word document."""
    doc = docx.Document(docx_path)
    text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
    return text

def extract_text_from_pptx(pptx_path):
    """Extracts text from a PowerPoint presentation."""
    presentation = pptx.Presentation(pptx_path)
    text = "\n".join(
        shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")
    )
    return text
